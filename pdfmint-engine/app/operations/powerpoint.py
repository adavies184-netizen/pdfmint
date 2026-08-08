from __future__ import annotations

import os
import re
import shutil
import subprocess
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches

from ..settings import JOB_TIMEOUT_SECONDS


def _remove_default_slide(prs: Presentation) -> None:
    if not prs.slides:
        return
    slide_id = prs.slides._sldIdLst[0]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[0]


def pdf_to_pptx(pdf_path: Path, output_path: Path) -> Path:
    document = fitz.open(pdf_path)

    try:
        if document.page_count < 1:
            raise RuntimeError("The PDF contains no pages.")

        presentation = Presentation()
        _remove_default_slide(presentation)

        first_page = document[0]
        aspect_ratio = first_page.rect.width / first_page.rect.height

        presentation.slide_height = Inches(7.5)
        presentation.slide_width = int(presentation.slide_height * aspect_ratio)

        blank_layout = presentation.slide_layouts[6]
        render_dir = output_path.parent / "ppt-pages"
        render_dir.mkdir(parents=True, exist_ok=True)

        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = render_dir / f"page-{page_index}.jpg"
            pixmap.save(str(image_path), jpg_quality=88)

            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        presentation.save(output_path)

    finally:
        document.close()

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise RuntimeError("PPTX conversion did not produce a valid file.")

    return output_path


def pptx_to_ppt(pptx_path: Path, output_dir: Path) -> Path:
    output_path = output_dir / f"{pptx_path.stem}.ppt"
    if output_path.exists():
        output_path.unlink()

    candidates = [
        "ppt",
        'ppt:"MS PowerPoint 97"',
        'ppt:"MS PowerPoint 95"',
    ]
    failures = []

    for convert_to in candidates:
        slug = re.sub(r"[^a-z0-9]+", "-", convert_to.lower()).strip("-")
        profile = output_dir / f"libreoffice-impress-{slug}"
        shutil.rmtree(profile, ignore_errors=True)
        profile.mkdir(parents=True, exist_ok=True)

        command = [
            "soffice",
            "--headless",
            "--invisible",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--nofirststartwizard",
            "--norestore",
            f"-env:UserInstallation=file://{profile}",
            "--convert-to",
            convert_to,
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ]

        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=JOB_TIMEOUT_SECONDS,
            check=False,
            env={
                **os.environ,
                "HOME": str(profile),
                "SAL_USE_VCLPLUGIN": "svp",
                "JAVA_TOOL_OPTIONS": "-Xms16m -Xmx64m",
            },
        )

        if (
            result.returncode == 0
            and output_path.exists()
            and output_path.stat().st_size > 0
        ):
            return output_path

        details = (result.stderr or result.stdout or "No PPT output was produced.").strip()
        failures.append(f"{convert_to}: exit={result.returncode}; {details[:500]}")

    raise RuntimeError(
        "LibreOffice could not export PPT using the available filters. "
        + " | ".join(failures)
    )


def pdf_to_ppt(pdf_path: Path, output_dir: Path, base_name: str) -> Path:
    pptx_path = output_dir / f"{base_name}.pptx"
    pdf_to_pptx(pdf_path, pptx_path)
    return pptx_to_ppt(pptx_path, output_dir)
